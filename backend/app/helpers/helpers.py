from pypdf import PdfReader
import faiss
from sentence_transformers import SentenceTransformer
from fastapi import UploadFile
from app.constant import index_path, metadata_path
from docx import Document
from pptx import Presentation
import numpy as np
import json
import re
import uuid
import os

# Function to extract text from a PDF file
# This function reads the PDF file and extracts text from specified pages
# If start_page and end_page are not specified, it reads the entire document
def extract_pdf_text(file_upload: UploadFile, start_page=0, end_page=None) -> str:
    # Make sure file pointer is at start
    file_upload.file.seek(0)

    pdf_reader = PdfReader(file_upload.file)
    total_pages = len(pdf_reader.pages)

    if end_page is None or end_page > total_pages:
        end_page = total_pages

    full_text = ""

    for page_num in range(start_page, end_page):
        page = pdf_reader.pages[page_num]
        text = page.extract_text()
        if text:
            full_text += text

    return full_text

# Function to extract text from a DOCX file
# This function reads the DOCX file and extracts text from all paragraphs
# It returns the text as a single string with paragraphs separated by newlines
def extract_docx_text(file_upload: UploadFile) -> str:
    file_upload.file.seek(0)
    doc = Document(file_upload.file)
    full_text = []

    # Extract paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text.strip())

    return '\n'.join(full_text)


def extract_pptx_text(file_upload: UploadFile) -> str:
    file_upload.file.seek(0)
    presentation = Presentation(file_upload.file)
    full_text = []

    # Extract text from each slide
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                full_text.append(shape.text.strip())

    return '\n'.join(full_text)

# Function to chunk text with context
# This function splits the text into chunks while preserving context
# It identifies headings and uses them to create a context stack
def chunk_text_with_context(text, chunk_size=5, overlap=2):
    lines = text.split('\n')
    chunks = []
    buffer = []
    context_stack = []

    def is_heading(line):
        return (
            bool(re.match(r"^[📌🧠🎯❓🧱⚙️🗃️📤🧾🧩🔄✅•]+\s+", line)) or  # Starts with emoji
            line.strip().endswith(':') or
            line.strip().endswith('?') or
            line.strip().isupper()
        )

    def update_context(line):
        if is_heading(line):
            context_stack.clear()
            context_stack.append(line.strip())

    for line in lines:
        line = line.strip()
        if not line:
            continue
        if is_heading(line):
            update_context(line)
            continue  # don't include heading itself in the content

        buffer.append((line, list(context_stack)))

    # Chunking logic with sliding window
    for i in range(0, len(buffer), chunk_size - overlap):
        chunk_lines = buffer[i:i + chunk_size]
        chunk_text = " ".join([line for line, _ in chunk_lines])
        chunk_context = " > ".join(chunk_lines[0][1]) if chunk_lines else "GENERAL"
        chunks.append({
            "content": chunk_text,
            "context": chunk_context
        })

    return chunks

# Function to embed text chunks
# This function takes the text chunks and embeds them using a provided embedder
def embed_text_chunks(chunks, model):
    texts = [chunk["content"] for chunk in chunks]
    return model.encode(texts, convert_to_numpy=True)

def build_faiss_index(embeddings, chunks):
    index = faiss.IndexFlatL2(embeddings.shape[1])
    index.add(np.array(embeddings, dtype=np.float32))
    
    # Store the original chunks for retrieval
    index.chunks = chunks
    return index

def build_faiss_index(embeddings):
    dim = embeddings.shape[1]
    index = faiss.IndexFlatL2(dim)
    index.add(embeddings)
    return index

def build_vector_dataset(chunk_texts, contexts, embeddings, source, topic="general"):
    assert len(chunk_texts) == len(contexts) == len(embeddings), \
        "Length mismatch: chunk_texts, contexts, and embeddings must be the same length."

    data = []
    for i in range(len(chunk_texts)):
        entry = {
            "id": str(uuid.uuid4()),
            "content": chunk_texts[i],
            "context": contexts[i],
            "topic": topic,
            "embedding": embeddings[i].tolist(),
            "source": source
        }
        data.append(entry)
    return data


def process_document_to_faiss(
    file_upload: UploadFile,
    start_page=1,
    end_page=None,
    topic="general",
    source=None,
    index=None,                
    vector_data=None           
):
    if source is None:
        source = os.path.basename(file_upload.filename)

    # Initialize vector_data if not provided
    if vector_data is None:
        vector_data = []

    print(f"Extracting text from: {source}")
    if file_upload.filename.endswith('.pdf'):
        text = extract_pdf_text(file_upload, start_page - 1, end_page)
    elif file_upload.filename.endswith('.docx'):
        text = extract_docx_text(file_upload)
    elif file_upload.filename.endswith('.pptx'):
        text = extract_pptx_text(file_upload)
    else:
        raise ValueError("Only PDF files are supported currently.")

    print("Chunking text...")
    chunks = chunk_text_with_context(text)

    print("Loading SentenceTransformer model...")
    embed_model = SentenceTransformer("all-MiniLM-L6-v2")

    print("Generating embeddings...")
    embeddings = embed_text_chunks(chunks, embed_model).astype("float32")

    chunk_texts = [chunk["content"] for chunk in chunks]
    contexts = [chunk["context"] for chunk in chunks]

    print("Building vector dataset entries...")
    new_vector_data = build_vector_dataset(chunk_texts, contexts, embeddings, source, topic=topic)

    print("Updating FAISS index...")
    if index is None:
        # Create a new index
        index = build_faiss_index(embeddings)
    else:
        # Check embedding dimension matches
        if embeddings.shape[1] != index.d:
            raise ValueError(f"Dimension mismatch! Index dim={index.d}, new vectors dim={embeddings.shape[1]}")
        index.add(embeddings)  # Append to existing index

    # Append new vector entries to dataset
    vector_data.extend(new_vector_data)

    save_faiss_index(index, index_path)
    save_metadata(vector_data, metadata_path)

    print(f"Added {len(new_vector_data)} new chunks from '{source}'")
    print(f"Total vectors in index: {index.ntotal}")
    return index, vector_data

def save_faiss_index(index, file_path):
    # cpu_index = faiss.index_gpu_to_cpu(index)  
    faiss.write_index(index, file_path)
    print(f"FAISS index saved to {file_path}")

def load_faiss_index(file_path):
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"FAISS index file not found: {file_path}")

    # res = faiss.StandardGpuResources()  # Tạo resource GPU
    cpu_index = faiss.read_index(file_path)  # Đọc index từ file (CPU)
    # gpu_index = faiss.index_cpu_to_gpu(res, 0, cpu_index)  # Chuyển sang GPU device 0
    print(f"FAISS index loaded from {file_path}")
    # return gpu_index
    return cpu_index

def load_metadata(file_path):
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Metadata file not found: {file_path}")
    with open(file_path, 'r', encoding='utf-8') as f:
        metadata = json.load(f)
    print(f"Metadata loaded from {file_path}")
    return metadata

def save_metadata(metadata, file_path):
    if not os.path.exists(os.path.dirname(file_path)):
        os.makedirs(os.path.dirname(file_path))
    with open(file_path, 'w', encoding='utf-8') as f:
        json.dump(metadata, f, ensure_ascii=False, indent=4)
    print(f"Metadata saved to {file_path}")